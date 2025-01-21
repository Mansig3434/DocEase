from docx import Document
import os
from django.core.files.storage import FileSystemStorage
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle
from reportlab.platypus import Paragraph, SimpleDocTemplate
from docx.shared import Pt,RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from pptx import Presentation


# merge document
def merge_doc_files(files):
    merged_document = Document()

    for file in files:
        doc = Document(file)
        for element in doc.element.body:
            merged_document.element.body.append(element)

    merged_path = os.path.join('media', 'merged_docs', 'merged_document.docx')

    os.makedirs(os.path.dirname(merged_path), exist_ok=True)

    merged_document.save(merged_path)

    return merged_path


# split document
def split_docx_by_sections(docx_path, fs: FileSystemStorage):
    document = Document(docx_path)
    split_docs = []
    section_counter = 1
    current_section = Document()
    for paragraph in document.paragraphs:
        current_section.add_paragraph(paragraph.text)
        if not paragraph.text.strip():
            section_filename = f'split_section_{section_counter}.docx'
            section_path = os.path.join(fs.location, section_filename)
            current_section.save(section_path)
            split_docs.append(section_filename)

            section_counter += 1
            current_section = Document()

    if current_section.paragraphs:
        section_filename = f'split_section_{section_counter}.docx'
        section_path = os.path.join(fs.location, section_filename)
        current_section.save(section_path)
        split_docs.append(section_filename)

    return split_docs


# doc_to_pdf

def convert_docx_to_pdf(doc_path_full, pdf_path):
    try:
        document = Document(doc_path_full)
        doc = SimpleDocTemplate(pdf_path, pagesize=letter)
        elements = []
        style = ParagraphStyle(
            name="Normal",
            fontName="Helvetica",
            fontSize=12,
            leading=14,
            spaceAfter=10
        )

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                elements.append(Paragraph(text, style))


        doc.build(elements)
    except Exception as e:
        print(f"Error during conversion: {e}")
        raise


#watermarking

def apply_watermark_to_doc(input_path, output_path, watermark_text):
    document = Document(input_path)

    footer = document.sections[0].footer
    footer_paragraph = footer.paragraphs[0]
    footer_paragraph.text = watermark_text
    run = footer_paragraph.runs[0]
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(192, 192, 192)
    run.font.bold = True
    run.font.name = "Arial"
    footer_paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

    header = document.sections[0].header
    header_paragraph = header.paragraphs[0]
    header_paragraph.text = watermark_text
    run = header_paragraph.runs[0]
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(192, 192, 192)
    run.font.bold = True
    run.font.name = "Arial"
    header_paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

    document.save(output_path)


#doc to ppt

def docx_to_pptx(input_path, output_path):
    document = Document(input_path)
    prs = Presentation()
    for para in document.paragraphs:
        if para.text.strip():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = "Slide Title"
            text_frame = content.text_frame
            p = text_frame.add_paragraph()
            p.text = para.text
    prs.save(output_path)

